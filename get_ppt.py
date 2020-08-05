import os
import time
import cx_Oracle
from selenium import webdriver
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.common.desired_capabilities import DesiredCapabilities
from pptx import Presentation
from pptx.util import Inches, Pt

#查询数据库函数
def sqlQuery(ymw):
	sql='''
	with t1 as(
	select t.yearnum,t.monthnum,t.weeknum,t.begin_day,t.end_day,t.month_end
	,t.yearnum||lpad(t.monthnum,2,0)||lpad(t.END_DAY,2,0) ymd 
	,case when t.monthnum=1 then 'Jan' 
	      when t.monthnum=2 then 'Feb'
	      when t.monthnum=3 then 'Mar'
	      when t.monthnum=4 then 'Apr'
	      when t.monthnum=5 then 'May'
	      when t.monthnum=6 then 'Jun'
	      when t.monthnum=7 then 'Jul'
	      when t.monthnum=8 then 'Aug'
	      when t.monthnum=9 then 'Sep'
	      when t.monthnum=10 then 'Oct'
	      when t.monthnum=11 then 'Nov'
	      when t.monthnum=12 then 'Dec' end month_en
	from dm_smart_week_define t
	)
	select t1.monthnum||'/'||t1.BEGIN_DAY||'~'||t1.monthnum||'/'||t1.END_DAY week_cn
	,t1.monthnum||'/1'||'~'||t1.monthnum||'/'||t1.END_DAY month_cn
	,t1.month_en||'.'||t1.BEGIN_DAY||'~'||t1.month_en||'.'||t1.END_DAY week_en
	,t1.month_en||'.1'||'~'||t1.month_en||'.'||t1.END_DAY month_en
	,t1.ymd,t1.monthnum||'月' month_id,'W'||t1.weeknum week_id,t1.month_end
	from t1
	where t1.yearnum||lpad(t1.monthnum,2,0)||'W'||t1.weeknum='{}'
	'''
	db_link= 'username/password@ip/server_name'
	con = cx_Oracle.connect(db_link)
	cur = con.cursor() 
	cur.execute(sql.format(ymw))
	rs = cur.fetchall()
	cur.close()
	con.close() 
	return rs

#截图函数
def TakeShot(driver, element, filename):
	item=WebDriverWait(driver,20).until(EC.presence_of_element_located((By.CSS_SELECTOR,element)))
	time.sleep(5)
	item.screenshot(filename+'.png')
	print(filename+'.png'+' 截图成功！')

#截图导出
def GetPhoto(month,week,month_end=0):
	#使用docker部署的selenium远程端执行
	driver = webdriver.Remote(command_executor="http://10.10.10.71:4444/wd/hub",desired_capabilities=DesiredCapabilities.CHROME)
	driver.implicitly_wait(10)        #隐性等待
	driver.get('http://telecom.thinktanksgmmd.com')
	driver.set_window_size(1920,1080) #浏览器分辨率设置
	#登录系统
	driver.find_element_by_id('username').send_keys('sgmm')
	driver.find_element_by_id('password').send_keys('1'+Keys.RETURN)

	#smart品牌人气跟踪页面
	driver.get('http://web.thinktanksgmmd.com/smart/mainMenu?pathUrl=#/smart-week-analysis/brand-popularity-track')
	time.sleep(4)
	#选择时间
	driver.find_element_by_css_selector('#root > div > div > div.page-content > div > div.toolbar.white > form > div:nth-child(1) > div > input').click()
	driver.find_element_by_link_text(month).click()
	driver.find_element_by_link_text(week).click()
	driver.find_element_by_css_selector('#root > div > div > div.page-content > div > div.toolbar.white > form > div:nth-child(1) > div > div > div > div.ui-button-container > button.ui-button.ui-confirm').click()
	#别克截图
	TakeShot(driver,'#root > div > div > div.page-content > div > div.p-10 > table.brand-popularity-track.reportInfo.for-sgm','./photo/buick_brand')
	#雪佛兰截图
	driver.find_element_by_link_text('雪佛兰').click()
	TakeShot(driver,'#root > div > div > div.page-content > div > div.p-10 > table.brand-popularity-track.reportInfo.for-sgm','./photo/Chevy_brand')
	#凯迪拉克截图
	driver.find_element_by_link_text('凯迪拉克APP').click()
	TakeShot(driver,'#root > div > div > div.page-content > div > div.p-10 > table.brand-popularity-track.reportInfo.for-sgm','./photo/Cadi_brand')

	#smart品牌人气汇总页面
	driver.get('http://web.thinktanksgmmd.com/smart/smart/smartMain.do#/smart-week-analysis/model-popularity-track')
	time.sleep(3)
	driver.find_element_by_css_selector('#root > div > div > div.page-content > div > div.toolbar.white > form > div:nth-child(1) > div > input').click()
	driver.find_element_by_link_text(month).click()
	driver.find_element_by_link_text(week).click()
	driver.find_element_by_css_selector('#root > div > div > div.page-content > div > div.toolbar.white > form > div:nth-child(1) > div > div > div > div.ui-button-container > button.ui-button.ui-confirm').click()
	driver.find_element_by_css_selector('#root > div > div > div.page-content > div > div:nth-child(2) > div.brand-option > div > ul > li:nth-child(2)').click()
	#选择别克车型
	driver.find_element_by_css_selector('#root > div > div > div.page-content > div > div.toolbar.white > form > div:nth-child(2) > div > a').click()
	driver.find_element_by_link_text('Velite 6(BEV)').click()
	driver.find_element_by_link_text('Velite 6(PHEV)').click()
	driver.find_element_by_css_selector('#root > div > div > div.page-content > div > div.toolbar.white > form > div:nth-child(2) > div > div > div > div.panel-body > div > div > div.tab-pane.active > div > div > div > table > tbody > tr:nth-child(1) > td:nth-child(2) > table > tbody > tr:nth-child(2) > td:nth-child(2) > table > tbody > tr:nth-child(2) > td:nth-child(2) > a:nth-child(1) > span:nth-child(3)').click()
	#拉动滚动条
	js='document.getElementsByClassName("scroller model-select-content")[0].scrollTop=100' 
	driver.execute_script(js)
	driver.find_element_by_link_text('Envision S').click()
	driver.find_element_by_link_text('GL8 Avenir').click()
	driver.find_element_by_link_text('GL8 ES').click()
	driver.find_element_by_link_text('GL8 25S/28T').click()
	driver.find_element_by_link_text('确定').click()
	#周度截图
	TakeShot(driver,'#root > div > div > div.page-content > div > div.p-10.reportInfo.brand-table-scroll > div > table','./photo/buick_model_week')
	driver.find_element_by_css_selector('#root > div > div > div.page-content > div > div:nth-child(2) > div.brand-option > div > ul > li:nth-child(3)').click()
	TakeShot(driver,'#root > div > div > div.page-content > div > div.p-10.reportInfo.brand-table-scroll > div > table','./photo/Chevy_model_week')
	driver.find_element_by_css_selector('#root > div > div > div.page-content > div > div:nth-child(2) > div.brand-option > div > ul > li:nth-child(4)').click()
	TakeShot(driver,'#root > div > div > div.page-content > div > div.p-10.reportInfo.brand-table-scroll > div > table','./photo/Cadi_model_week')
	#月末周需取总量
	if month_end==1:
		driver.find_element_by_link_text('总量').click()
	else:
		pass
	#月度截图
	driver.find_element_by_link_text('月累计').click()
	TakeShot(driver,'#root > div > div > div.page-content > div > div.p-10.reportInfo.brand-table-scroll > div > table','./photo/Cadi_model_month')
	driver.find_element_by_css_selector('#root > div > div > div.page-content > div > div:nth-child(2) > div.brand-option > div > ul > li:nth-child(3)').click()
	TakeShot(driver,'#root > div > div > div.page-content > div > div.p-10.reportInfo.brand-table-scroll > div > table','./photo/Chevy_model_month')
	driver.find_element_by_css_selector('#root > div > div > div.page-content > div > div:nth-child(2) > div.brand-option > div > ul > li:nth-child(2)').click()
	TakeShot(driver,'#root > div > div > div.page-content > div > div.p-10.reportInfo.brand-table-scroll > div > table','./photo/buick_model_month')
	#关闭浏览器
	driver.close()

#生成PPT函数
def add_slide(text,photo):
	if text[0] in (5,7):
		#title添加新文本
		# 插入幻灯片
		blank_slide = prs.slide_layouts[5] #使用第7个母版格式
		slide = prs.slides.add_slide(blank_slide)
		#title
		title_shape = slide.shapes.title 
		title_shape.text = text[1]
		#title添加新文本
		new_paragraph =title_shape.text_frame.add_paragraph()
		new_paragraph.text = text[2]
		new_paragraph.font.size = Pt(16)
		# 添加图片
		slide.shapes.add_picture(image_file=photo[0],left=Inches(0.5),top=Inches(1.4),width=Inches(9),height=Inches(2.8))
		slide.shapes.add_picture(image_file=photo[1],left=Inches(0.5),top=Inches(4.5),width=Inches(9),height=Inches(2.8))
	elif text[0]==2:
		# 插入幻灯片
		blank_slide = prs.slide_layouts[3] #使用第6个母版格式
		slide = prs.slides.add_slide(blank_slide)
		#title
		title_shape = slide.shapes.title 
		title_shape.text = text[1]
		#title添加新文本
		new_paragraph =title_shape.text_frame.add_paragraph()
		new_paragraph.text = text[2]
		new_paragraph.font.size = Pt(16)
		# 添加图片
		slide.shapes.add_picture(image_file=photo,left=Inches(0),top=Inches(1.5),width=Inches(10),height=Inches(5.5))
	elif text[0]==3:
		# 插入幻灯片
		blank_slide = prs.slide_layouts[4] #使用第6个母版格式
		slide = prs.slides.add_slide(blank_slide)
		#title
		title_shape = slide.shapes.title 
		title_shape.text = text[1]
		#title添加新文本
		new_paragraph =title_shape.text_frame.add_paragraph()
		new_paragraph.text = text[2]
		new_paragraph.font.size = Pt(16)
		# 添加图片
		slide.shapes.add_picture(image_file=photo,left=Inches(0),top=Inches(1.5),width=Inches(10),height=Inches(5.5))
	else:
		# 插入幻灯片
		blank_slide = prs.slide_layouts[2] #使用第6个母版格式
		slide = prs.slides.add_slide(blank_slide)
		#title
		title_shape = slide.shapes.title 
		title_shape.text = text[1]
		#title添加新文本
		new_paragraph =title_shape.text_frame.add_paragraph()
		new_paragraph.text = text[2]
		new_paragraph.font.size = Pt(16)
		# 添加图片
		slide.shapes.add_picture(image_file=photo,left=Inches(0),top=Inches(1.5),width=Inches(10),height=Inches(5.5))

if __name__ == "__main__":
	try:
		#获取Jenkins时间
		ymw=os.environ["smart_ymw"]
		#查询数据库信息
		timedata=sqlQuery(ymw.upper())
		#导出图片
		GetPhoto(timedata[0][5],timedata[0][6],timedata[0][7])

		#使用PPT模板
		prs = Presentation('./template/smart_temple.pptx')
		#标题名称
		text=[]
		text.append([1,'别克市场人气跟踪（{}）'.format(timedata[0][0]),'Buick Showroom Tracking ({})'.format(timedata[0][2])])
		text.append([2,'别克分车型市场人气跟踪（ {} ）'.format(timedata[0][0]),'Buick Showroom Tracking by Models ({})'.format(timedata[0][2])])
		text.append([3,'别克分车型市场人气跟踪（ {} ）'.format(timedata[0][1]),'Buick Showroom Tracking by Models ({})'.format(timedata[0][3])])
		text.append([4,'雪佛兰市场人气跟踪（ {} ）'.format(timedata[0][0]),'Chevy Showroom Tracking ({})'.format(timedata[0][2])])
		text.append([5,'雪佛兰分车型市场人气跟踪（ {} ）'.format(timedata[0][0]),'Chevrolet Showroom Tracking by Models ({})'.format(timedata[0][2])])
		text.append([6,'凯迪拉克市场人气跟踪（ {} ）'.format(timedata[0][0]),'Cadi Showroom Tracking ({})'.format(timedata[0][2])])
		text.append([7,'凯迪拉克分车型市场人气跟踪（ {} ）'.format(timedata[0][0]),'Cadillac Showroom Tracking by Models ({})'.format(timedata[0][2])])
		#图片名称
		photo=[]
		photo.append('./photo/buick_brand.png')
		photo.append('./photo/buick_model_week.png')
		photo.append('./photo/buick_model_month.png')
		photo.append('./photo/Chevy_brand.png')
		photo.append(['./photo/Chevy_model_week.png','./photo/Chevy_model_month.png'])
		photo.append('./photo/Cadi_brand.png')
		photo.append(['./photo/Cadi_model_week.png','./photo/Cadi_model_month.png'])
		#生成PPT
		for i in range(len(text)):
			add_slide(text[i],photo[i])
		#保存ppt
		prs.save('./result/{}-SMART报表功能.pptx'.format(timedata[0][4]))
		print('自动化报告【{}-SMART报表功能.pptx】已生成！'.format(timedata[0][4]))
		
	except Exception as e:
		print(e)
	

